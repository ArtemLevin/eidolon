from __future__ import annotations
from typing import List
from io import BytesIO
from PIL import Image
from pptx import Presentation

def ingest_pptx(pptx_path: str, max_slides: int = 20) -> tuple[str, List[bytes]]:
    prs = Presentation(pptx_path)
    texts = []
    image_bytes: List[bytes] = []
    for i, slide in enumerate(prs.slides):
        if i >= max_slides:
            break
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if shape.shape_type == 13 and not image_bytes:
                try:
                    img = Image.open(BytesIO(shape.image.blob))
                    buf = BytesIO(); img.save(buf, format="PNG")
                    image_bytes.append(buf.getvalue())
                except Exception:
                    pass
    return "\n".join(texts).strip(), image_bytes[:1]
